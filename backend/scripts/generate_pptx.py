import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define corporate colors (TEI style - assuming Blue/Gold/White based on description)
    TEI_BLUE = RGBColor(0, 51, 102)
    TEI_GOLD = RGBColor(218, 165, 32)
    TEI_WHITE = RGBColor(255, 255, 255)
    TEI_GREY = RGBColor(100, 100, 100)

    def add_slide(title_text, content_text_list, layout_index=1):
        """Helper to add a slide with title and bullet points."""
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = TEI_BLUE
        title.text_frame.paragraphs[0].font.bold = True

        # Content
        if content_text_list:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = content_text_list[0]
            
            for item in content_text_list[1:]:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0

    # --- SLIDE 1: TITLE ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "CENTRO COMERCIAL TEI"
    subtitle.text = "Tu Futuro Financiero Comienza Aquí\n\nCompra, Ahorra y Gana en Dólares"
    
    title.text_frame.paragraphs[0].font.color.rgb = TEI_BLUE
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text_frame.paragraphs[0].font.color.rgb = TEI_GOLD

    # --- SLIDE 2: PROBLEM VS SOLUTION ---
    add_slide(
        "El Problema vs. La Solución",
        [
            "¿Dependes de una sola fuente de ingresos?",
            "La inflación y las deudas afectan tu tranquilidad.",
            "",
            "LA SOLUCIÓN TEI:",
            "✅ Ecosistema Híbrido: E-commerce + Network Marketing.",
            "✅ 4 Formas Simultáneas de Ganar.",
            "✅ Ingresos en Dólares y Cripto.",
            "✅ Sin dejar tu ocupación actual."
        ]
    )

    # --- SLIDE 3: NUESTRO ECOSISTEMA ---
    add_slide(
        "Nuestro Ecosistema Integral",
        [
            "1. Tienda Virtual: Productos físicos y digitales de alta demanda.",
            "2. Plan Binario Global: Ingreso por posicionamiento rápido.",
            "3. Plan Binario Millonario: Para constructores de equipos y liderazgo.",
            "4. Bonos de Honor: Reparto de utilidades globales.",
            "",
            "Base sólida de comercio real con un motor de comisiones explosivo."
        ]
    )

    # --- SLIDE 4: PLAN BINARIO GLOBAL ---
    add_slide(
        "Plan Binario Global (La Joya de la Corona)",
        [
            "🚀 POSICIONAMIENTO AUTOMÁTICO",
            "Los usuarios se colocan en una red 2x2 por orden de llegada mundial.",
            "",
            "🌊 DERRAME MUNDIAL (SPILLOVER)",
            "Recibe beneficios del crecimiento global de la empresa.",
            "",
            "⏳ PERIODO DE GRACIA",
            "4 Meses para activarte sin perder tu posición."
        ]
    )

    # --- SLIDE 5: GANANCIAS BINARIO GLOBAL ---
    add_slide(
        "Potencial de Ganancias: Binario Global",
        [
            "¡GANA POR AMBAS LÍNEAS (Izquierda y Derecha)!",
            "Pagos por niveles impares hasta el NIVEL 21:",
            "",
            "🔹 Niveles 3 al 13: $0.50 USD por persona.",
            "🔹 Niveles 15 al 21: $1.00 USD por persona.",
            "",
            "EJEMPLO DE PODER:",
            "• En una red 2x2 completa, el nivel 21 tiene más de 2 millones de posiciones.",
            "• Incluso llenando una fracción, el ingreso es masivo.",
            "• No importa si están en tu pierna fuerte o débil: ¡Cobras por todos!",
            "",
            "Requisito: Activación con compra mínima."
        ]
    )

    # --- SLIDE 6: PLAN BINARIO MILLONARIO ---
    add_slide(
        "Plan Binario Millonario",
        [
            "Diseñado para Líderes y Constructores.",
            "",
            "📊 Puntos Volumen (PV)",
            "Cada producto suma puntos.",
            "",
            "⚖️ Equipo Menor",
            "Cobra un porcentaje (ej. 10%) del volumen de tu pierna menor.",
            "",
            "💰 Profundidad hasta Nivel 27",
            "Pagos Rápidos. Cortes diarios o semanales."
        ]
    )

    # --- SLIDE 7: UNILEVEL Y MATRIZ ---
    add_slide(
        "Ingreso Residual: Unilevel y Matriz",
        [
            "La verdadera libertad financiera.",
            "",
            "🏢 PLAN UNILEVEL",
            "Gana un porcentaje de las compras directas de tus referidos.",
            "Hasta 7 niveles de profundidad.",
            "",
            "🕸️ MATRIZ FORZADA CERRADA 3x3",
            "Estructura compacta y poderosa:",
            "• Nivel 1: 3 Personas",
            "• Nivel 2: 9 Personas",
            "• Total: 12 Personas para completar ciclo.",
            "Ingresos residuales garantizados por la compra de tu primer paquete"
        ]
    )

    # --- SLIDE 8: CARRERA DE RANGOS ---
    add_slide(
        "Carrera de Honor y Liderazgo",
        [
            "Reconocemos tu esfuerzo y resultados.",
            "",
            "🏆 RANGOS",
            "Plata -> Oro -> Diamante -> Embajador.",
            "",
            "🌍 POOL GLOBAL (Rangos de Honor)",
            "La empresa reparte un % de las ventas mundiales entre los líderes calificados.",
            "¡Conviértete en socio de la compañía!"
        ]
    )

    # --- SLIDE 9: CÓMO INICIAR ---
    add_slide(
        "¿Cómo Iniciar Hoy?",
        [
            "1️⃣ PRE-REGÍSTRATE GRATIS",
            "Asegura tu lugar en el Binario Global AHORA.",
            "",
            "2️⃣ ACTÍVATE",
            "Compra tu paquete de inicio o productos en la tienda.",
            "",
            "3️⃣ COMPARTE",
            "Usa tu enlace de referido y nuestro sistema automático.",
            "",
            "¡El tiempo es dinero! Posiciónate antes que el resto."
        ]
    )

    # --- SLIDE 10: CIERRE ---
    slide_layout = prs.slide_layouts[0] # Title Slide layout for closing
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "TU MOMENTO ES AHORA"
    subtitle.text = "Contacta a la persona que te invitó y asegura tu posición.\n\nCentro Comercial TEI"
    
    title.text_frame.paragraphs[0].font.color.rgb = TEI_BLUE

    # Save
    output_path = os.path.join("docs", "Presentacion_Negocio_TEI.pptx")
    prs.save(output_path)
    print(f"✅ Presentación creada exitosamente en: {output_path}")

if __name__ == "__main__":
    create_presentation()
